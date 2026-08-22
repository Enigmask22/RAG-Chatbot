"""Sinh 7 file fixture cho `W3-01`/`W3-02`, một file cho mỗi định dạng.

Chạy: `make loader-fixtures`

**Vì sao sinh chứ không đi tìm file thật.** Quy tắc corpus của dự án (`W1-02`)
là mọi thứ đi vào repo phải công khai và cho phép redistribute — repo public +
demo public + máy GPU thuê là ba kênh công bố. Một PDF hai cột "mượn tạm" ở đâu
đó vi phạm đúng quy tắc ấy. Nhưng lý do mạnh hơn là **kiểm chứng**: với file
thật tôi chỉ có thể nhìn output rồi gật đầu, còn với file tự sinh thì thứ tự
đọc đúng là **biết trước theo cách dựng**, nên test khẳng định được chứ không
phải cảm nhận được.

**Chỗ then chốt: PDF ở đây cố ý ghi content stream SAI thứ tự đọc.** Các dòng
được ghi xen kẽ trái–phải–trái–phải, tức đúng thứ tự mà một bộ trích text chạy
theo content stream sẽ trả ra. PDF thật hay hỏng đúng kiểu này, vì thứ tự trong
content stream là thứ tự **trình bày**, không phải thứ tự **đọc**; không có gì
trong đặc tả PDF bắt hai thứ đó trùng nhau. Nên fixture này phân biệt được hai
thứ mà nhìn output không phân biệt nổi:

* bộ trích theo content stream → `Alpha …, Beta …, Alpha …, Beta …`
* bộ có phân tích bố cục       → trọn cột `Alpha` rồi mới tới cột `Beta`

Không dùng thư viện sinh PDF nào: PDF tối giản với font base-14 chỉ là văn bản
có bảng offset, viết tay ~40 dòng, và đổi lại là fixture không phụ thuộc vào
version của một thư viện thứ ba mà tôi lại đang muốn đo tính ổn định.

**Ba file OOXML được "đóng băng" zip.** `.docx`/`.pptx`/`.xlsx` là zip, mà zip
ghi mtime từng entry và OOXML ghi thêm ngày tạo/sửa vào `docProps`. Chạy lại
script hai lần sẽ ra hai chuỗi byte khác nhau, tức `git status` bẩn mỗi lần và
sha256 trong test vô nghĩa. `_freeze_zip` ghi lại toàn bộ archive với timestamp
cố định 1980-01-01 để `make loader-fixtures` là **idempotent**.
"""

from __future__ import annotations

import argparse
import datetime as dt
import hashlib
import io
import re
import textwrap
import zipfile
from pathlib import Path

__all__ = ["FIXTURE_DIR", "write_fixtures"]

FIXTURE_DIR = Path(__file__).resolve().parents[1] / "tests" / "fixtures" / "loaders"

_FROZEN_TIME = dt.datetime(1980, 1, 1, 0, 0, 0, tzinfo=dt.UTC)
_FROZEN_ZIP_TIME = (1980, 1, 1, 0, 0, 0)

# Cột trái và cột phải của `two-column.pdf`. Thứ tự đọc đúng là hết ALPHA rồi
# mới tới BETA; content stream thì ghi xen kẽ.
#
# ⚠️ Hai cột này CỐ Ý là văn xuôi có độ dài dòng so le, không phải mấy dòng ngắn
# giống hệt nhau. Bản fixture đầu tiên của tôi đúng kiểu "sạch" ấy và nó **hỏng**:
# model bố cục của docling được huấn luyện trên ảnh trang thật, còn một trang
# gồm vài dòng ngắn bằng nhau cách đều nhau thì nằm ngoài phân bố đó, và kết quả
# nhảy không đơn điệu theo số dòng — 4 dòng ✗, 12 ✓, 24 ✗ (cột phải bị **phân
# loại thành bảng**), 40 ✓. Với cột văn xuôi so le thì 4/4 biến thể đều đúng.
# Nếu tôi giữ fixture "sạch" thì test đo cái generator của tôi chạm vào biên
# quyết định của model ở đâu, chứ không đo docling có đọc được hai cột không.
_PDF_TITLE = "Two Column Reading Order Probe"
_PDF_LEFT_PROSE = (
    "Alpha. The left column carries the first half of the argument and it runs "
    "for several sentences so that the block of text on this side of the page "
    "has the ragged silhouette that a real typeset paragraph has rather than a "
    "stack of identical lines. It mentions figures such as 5,05 percent and "
    "7,09 percent, and it refers to the year 2024 more than once, because "
    "prose in a report tends to."
)
_PDF_RIGHT_PROSE = (
    "Beta. The right column continues on the other side of the gutter and it "
    "is written to a different length so the two blocks do not line up neatly "
    "at the bottom edge. It discusses inflation of 3,25 percent, exports worth "
    "355,5 billion dollars, and a handful of other numbers that would appear "
    "in a macroeconomic bulletin of this kind."
)
_PDF_WRAP_CHARS = 34
_PDF_FONT_SIZE = 9
_PDF_LINE_STEP = 12

# Cùng một bảng ở cả bốn định dạng có bảng — để so được "bảng có giữ cấu trúc
# không" giữa các backend mà không phải nhớ mỗi file một nội dung.
_TABLE_HEADER = ("Chỉ tiêu", "2023", "2024")
_TABLE_ROWS = (
    ("Tăng trưởng GDP (%)", "5,05", "7,09"),
    ("Lạm phát CPI (%)", "3,25", "3,63"),
    ("Xuất khẩu (tỷ USD)", "355,5", "405,5"),
)


# --------------------------------------------------------------------------
# PDF — viết tay, không thư viện
# --------------------------------------------------------------------------
def _escape_pdf_text(text: str) -> bytes:
    """Escape theo cú pháp chuỗi literal của PDF, mã hoá WinAnsi.

    Chỉ ASCII đi vào PDF: font base-14 dùng WinAnsiEncoding, không có dấu tiếng
    Việt. Fixture này đo bố cục chứ không đo ngôn ngữ, còn phần tiếng Việt đã
    có ở 5 định dạng kia.
    """
    escaped = text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
    return escaped.encode("ascii")


def _text_object(x: float, y: float, size: int, text: str) -> bytes:
    return (
        b"BT\n/F1 "
        + str(size).encode("ascii")
        + b" Tf\n1 0 0 1 "
        + f"{x:g} {y:g}".encode("ascii")
        + b" Tm\n("
        + _escape_pdf_text(text)
        + b") Tj\nET\n"
    )


def _pdf_content_stream() -> bytes:
    """Content stream với các dòng **xen kẽ hai cột** — sai thứ tự đọc có chủ đích."""
    left = textwrap.wrap(_PDF_LEFT_PROSE, _PDF_WRAP_CHARS)
    right = textwrap.wrap(_PDF_RIGHT_PROSE, _PDF_WRAP_CHARS)
    parts = [_text_object(72, 752, 14, _PDF_TITLE)]
    for row in range(max(len(left), len(right))):
        y = 720 - row * _PDF_LINE_STEP
        if row < len(left):
            parts.append(_text_object(72, y, _PDF_FONT_SIZE, left[row]))
        if row < len(right):
            parts.append(_text_object(330, y, _PDF_FONT_SIZE, right[row]))
    return b"".join(parts)


def _assemble_pdf(objects: list[bytes]) -> bytes:
    """Ghép các object thành file PDF hợp lệ, tự tính bảng xref."""
    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode("ascii") + body + b"\nendobj\n"
    xref_offset = len(out)
    size = len(objects) + 1
    out += f"xref\n0 {size}\n".encode("ascii")
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode("ascii")
    out += (f"trailer\n<< /Size {size} /Root 1 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n").encode(
        "ascii"
    )
    return bytes(out)


def build_pdf() -> bytes:
    stream = _pdf_content_stream()
    return _assemble_pdf(
        [
            b"<< /Type /Catalog /Pages 2 0 R >>",
            b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
            b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica /Encoding /WinAnsiEncoding >>",
            b"<< /Length "
            + str(len(stream)).encode("ascii")
            + b" >>\nstream\n"
            + stream
            + b"endstream",
        ]
    )


# --------------------------------------------------------------------------
# PDF scan — ảnh một trang, KHÔNG có text layer (`W3-02`)
# --------------------------------------------------------------------------
# Hai đoạn cố ý khác ngôn ngữ. Máy OCR mặc định của docling là RapidOCR với model
# PP-OCRv6 + `ppocrv6_dict.txt` — bộ chữ Trung/Anh. Dự án này phục vụ corpus tiếng
# Việt, nên "OCR có đọc được tiếng Việt có dấu không" là câu phải ĐO chứ không
# phải câu để giả định. Đặt hai đoạn cạnh nhau trong cùng một ảnh thì so được
# trực tiếp, cùng độ phân giải, cùng font, cùng một lần chạy.
_SCAN_EN_LINES = (
    "Vietnam Macroeconomic Update",
    "Growth reached 7.09 percent in 2024, up",
    "from 5.05 percent a year earlier. Exports",
    "rose to 405.5 billion dollars while the",
    "consumer price index settled at 3.63.",
)
_SCAN_VI_LINES = (
    "Cap nhat kinh te vi mo Viet Nam",
    "Tăng trưởng đạt 7,09 phần trăm năm 2024,",
    "cao hơn mức 5,05 phần trăm của năm trước.",
    "Xuất khẩu tăng lên 405,5 tỷ đô la trong khi",
    "chỉ số giá tiêu dùng dừng ở mức 3,63.",
)
_SCAN_DPI = 150
_SCAN_PAGE_PX = (int(8.5 * _SCAN_DPI), int(11 * _SCAN_DPI))
# ⚠️ **Cả hai** trường, không chỉ `/CreationDate`. Vá một cái rồi tưởng xong là
# đúng lỗi đã mắc với `docProps` của openpyxl ở trên: hai lần chạy cách nhau
# dưới một giây thì trùng, cách nhau hơn một giây thì lệch **đúng một byte**.
_PDF_DATE_RE = re.compile(rb"/(?:Creation|Mod)Date\s*\(D:[^)]*\)")


def build_scanned_pdf() -> bytes:
    """Một trang giấy trắng có chữ, lưu thành PDF **ảnh** — không text layer.

    Pillow ghi `/CreationDate` bằng giờ hiện tại nên phải ép lại, cùng lý do với
    `docProps` của OOXML ở `_freeze_zip`. Font dùng `ImageFont.load_default(size=)`
    (Aileron, đóng gói sẵn trong Pillow) chứ không dùng font hệ điều hành: fixture
    sinh trên máy khác phải ra cùng chuỗi byte.
    """
    from PIL import Image, ImageDraw, ImageFont

    page = Image.new("L", _SCAN_PAGE_PX, color=255)
    draw = ImageDraw.Draw(page)
    title_font = ImageFont.load_default(size=34)
    body_font = ImageFont.load_default(size=26)

    y = 120
    for block in (_SCAN_EN_LINES, _SCAN_VI_LINES):
        for index, line in enumerate(block):
            draw.text((110, y), line, fill=0, font=title_font if index == 0 else body_font)
            y += 52 if index == 0 else 44
        y += 60

    buffer = io.BytesIO()
    # `resolution` để pypdfium2 báo đúng kích thước trang (8,5 × 11 inch).
    page.convert("1").save(buffer, format="PDF", resolution=_SCAN_DPI)
    return _PDF_DATE_RE.sub(
        lambda m: m.group(0)[: m.group(0).index(b"(")] + b"(D:19800101000000Z)",
        buffer.getvalue(),
    )


# --------------------------------------------------------------------------
# Markdown / HTML — văn bản thuần, không phụ thuộc gì
# --------------------------------------------------------------------------
def build_md() -> bytes:
    table = ["| " + " | ".join(_TABLE_HEADER) + " |", "| --- | --- | --- |"]
    table += ["| " + " | ".join(row) + " |" for row in _TABLE_ROWS]
    lines = [
        "# Chương I — Tổng quan kinh tế",
        "",
        "Đoạn mở đầu của chương, nằm ngay dưới heading cấp 1.",
        "",
        "## Điều 1. Phạm vi điều chỉnh",
        "",
        "Đoạn thuộc heading cấp 2.",
        "",
        "### Khoản 1. Chỉ tiêu vĩ mô",
        "",
        *table,
        "",
        "Đoạn kết thuộc heading cấp 3.",
        "",
    ]
    return "\n".join(lines).encode("utf-8")


def build_html() -> bytes:
    rows = "\n".join(
        "      <tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>" for row in _TABLE_ROWS
    )
    header = "".join(f"<th>{cell}</th>" for cell in _TABLE_HEADER)
    html = f"""<!DOCTYPE html>
<html lang="vi">
  <head><meta charset="utf-8"><title>Chương I — Tổng quan kinh tế</title></head>
  <body>
    <h1>Chương I — Tổng quan kinh tế</h1>
    <p>Đoạn mở đầu của chương, nằm ngay dưới heading cấp 1.</p>
    <h2>Điều 1. Phạm vi điều chỉnh</h2>
    <p>Đoạn thuộc heading cấp 2.</p>
    <h3>Khoản 1. Chỉ tiêu vĩ mô</h3>
    <table>
      <tr>{header}</tr>
{rows}
    </table>
    <p>Đoạn kết thuộc heading cấp 3.</p>
  </body>
</html>
"""
    return html.encode("utf-8")


# --------------------------------------------------------------------------
# OOXML — docx / pptx / xlsx
# --------------------------------------------------------------------------
_CORE_PROPS = "docProps/core.xml"
_DCTERMS_RE = re.compile(rb"(<dcterms:(?:created|modified)[^>]*>)[^<]*(</dcterms:)")


def _freeze_core_props(payload: bytes) -> bytes:
    """Ép `dcterms:created`/`modified` về mốc cố định.

    Đặt `core_properties.modified` trước khi save là **không đủ**: openpyxl ghi
    đè bằng giờ hiện tại ngay trong lúc serialize, nên hai lần chạy cách nhau
    một giây ra hai file khác nhau. Đo được: 6 tiến trình cho 3 chuỗi byte, khác
    nhau đúng ở trường này. Phải sửa **sau** khi save.
    """
    stamp = _FROZEN_TIME.strftime("%Y-%m-%dT%H:%M:%SZ").encode("ascii")
    return _DCTERMS_RE.sub(rb"\g<1>" + stamp + rb"\g<2>", payload)


def _freeze_zip(payload: bytes) -> bytes:
    """Ghi lại archive với timestamp cố định để việc sinh fixture là idempotent."""
    source = zipfile.ZipFile(io.BytesIO(payload))
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as target:
        for name in sorted(source.namelist()):
            body = source.read(name)
            if name == _CORE_PROPS:
                body = _freeze_core_props(body)
            info = zipfile.ZipInfo(name, date_time=_FROZEN_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            target.writestr(info, body)
    source.close()
    return buffer.getvalue()


def build_docx() -> bytes:
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_heading("Chương I — Tổng quan kinh tế", level=1)
    doc.add_paragraph("Đoạn mở đầu của chương, nằm ngay dưới heading cấp 1.")
    doc.add_heading("Điều 1. Phạm vi điều chỉnh", level=2)
    doc.add_paragraph("Đoạn thuộc heading cấp 2.")
    doc.add_heading("Khoản 1. Chỉ tiêu vĩ mô", level=3)
    table = doc.add_table(rows=1 + len(_TABLE_ROWS), cols=len(_TABLE_HEADER))
    table.style = "Table Grid"
    for column, value in enumerate(_TABLE_HEADER):
        table.cell(0, column).text = value
    for row_index, row in enumerate(_TABLE_ROWS, start=1):
        for column, value in enumerate(row):
            table.cell(row_index, column).text = value
    doc.add_paragraph("Đoạn kết thuộc heading cấp 3.")

    props = doc.core_properties
    props.created = _FROZEN_TIME
    props.modified = _FROZEN_TIME
    props.title = "Chương I — Tổng quan kinh tế"
    props.author = "rag-platform fixtures"
    props.last_modified_by = "rag-platform fixtures"
    props.revision = 1

    buffer = io.BytesIO()
    doc.save(buffer)
    return _freeze_zip(buffer.getvalue())


def build_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    title_slide = deck.slides.add_slide(deck.slide_layouts[0])
    title_slide.shapes.title.text = "Chương I — Tổng quan kinh tế"
    title_slide.placeholders[1].text = "Đoạn mở đầu của chương, nằm ngay dưới heading cấp 1."

    body_slide = deck.slides.add_slide(deck.slide_layouts[1])
    body_slide.shapes.title.text = "Điều 1. Phạm vi điều chỉnh"
    frame = body_slide.placeholders[1].text_frame
    frame.text = "Khoản 1. Chỉ tiêu vĩ mô"
    for row in _TABLE_ROWS:
        paragraph = frame.add_paragraph()
        paragraph.text = f"{row[0]}: {row[1]} → {row[2]}"

    table_slide = deck.slides.add_slide(deck.slide_layouts[5])
    table_slide.shapes.title.text = "Khoản 2. Bảng chỉ tiêu"
    shape = table_slide.shapes.add_table(
        1 + len(_TABLE_ROWS), len(_TABLE_HEADER), Inches(0.5), Inches(1.8), Inches(9), Inches(2.4)
    )
    grid = shape.table
    for column, value in enumerate(_TABLE_HEADER):
        grid.cell(0, column).text = value
    for row_index, row in enumerate(_TABLE_ROWS, start=1):
        for column, value in enumerate(row):
            grid.cell(row_index, column).text = value

    props = deck.core_properties
    props.created = _FROZEN_TIME
    props.modified = _FROZEN_TIME
    props.title = "Chương I — Tổng quan kinh tế"
    props.author = "rag-platform fixtures"
    props.last_modified_by = "rag-platform fixtures"
    props.revision = 1

    buffer = io.BytesIO()
    deck.save(buffer)
    return _freeze_zip(buffer.getvalue())


def build_xlsx() -> bytes:
    from openpyxl import Workbook

    book = Workbook()
    sheet = book.active
    sheet.title = "Chỉ tiêu vĩ mô"
    sheet.append(list(_TABLE_HEADER))
    for row in _TABLE_ROWS:
        sheet.append(list(row))

    book.properties.created = _FROZEN_TIME.replace(tzinfo=None)
    book.properties.modified = _FROZEN_TIME.replace(tzinfo=None)
    book.properties.title = "Chỉ tiêu vĩ mô"
    book.properties.creator = "rag-platform fixtures"
    book.properties.lastModifiedBy = "rag-platform fixtures"

    buffer = io.BytesIO()
    book.save(buffer)
    return _freeze_zip(buffer.getvalue())


_BUILDERS = {
    "two-column.pdf": build_pdf,
    "scanned-page.pdf": build_scanned_pdf,
    "chuong-i.docx": build_docx,
    "chuong-i.pptx": build_pptx,
    "chi-tieu.xlsx": build_xlsx,
    "chuong-i.html": build_html,
    "chuong-i.md": build_md,
}


def write_fixtures(directory: Path = FIXTURE_DIR) -> dict[str, str]:
    """Ghi cả 6 fixture, trả `{tên file: sha256}`."""
    directory.mkdir(parents=True, exist_ok=True)
    digests: dict[str, str] = {}
    for name, builder in _BUILDERS.items():
        payload = builder()
        (directory / name).write_bytes(payload)
        digests[name] = hashlib.sha256(payload).hexdigest()
    return digests


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--out", type=Path, default=FIXTURE_DIR)
    args = parser.parse_args(argv)

    digests = write_fixtures(args.out)
    width = max(len(name) for name in digests)
    for name, digest in digests.items():
        size = (args.out / name).stat().st_size
        print(f"{name:<{width}}  {size:>7,d} B  {digest[:16]}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
